"""
Format handlers for non-PDF/DOCX documents.

Supports: Excel, Email (EML/MSG), PowerPoint, and fallback for other formats.
"""

import json
from pathlib import Path
from typing import Optional, Dict
from citation_types import ConversionResult, DocumentType


class FormatHandler:
    """Route documents to appropriate converters based on file extension."""

    SUPPORTED_FORMATS = {
        '.pdf': 'docling',
        '.docx': 'docling',
        '.xlsx': 'excel_parser',
        '.xls': 'excel_parser',
        '.pptx': 'powerpoint_parser',
        '.eml': 'email_parser',
        '.msg': 'email_parser',
        '.txt': 'text_parser',
        '.md': 'text_parser',
    }

    def __init__(self, output_dir: str):
        """
        Initialize format handler.

        Args:
            output_dir: Directory to save converted files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def convert(self, input_path: Path) -> ConversionResult:
        """
        Convert document using appropriate format handler.

        Args:
            input_path: Path to input document

        Returns:
            ConversionResult
        """
        ext = input_path.suffix.lower()
        converter = self.SUPPORTED_FORMATS.get(ext, 'textract_fallback')

        # Call appropriate converter method
        if hasattr(self, converter):
            return getattr(self, converter)(input_path)
        else:
            return self.textract_fallback(input_path)

    def excel_parser(self, input_path: Path) -> ConversionResult:
        """
        Convert Excel files to markdown.

        Strategy: Extract text from each sheet, preserve table structure.
        """
        try:
            import openpyxl
        except ImportError:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=["openpyxl not installed. Run: pip install openpyxl"]
            )

        try:
            workbook = openpyxl.load_workbook(input_path, data_only=True)
            md_lines = [f"# {input_path.stem}\n"]

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                md_lines.append(f"\n## Sheet: {sheet_name}\n")

                # Extract table data
                for row in sheet.iter_rows(values_only=True):
                    # Skip empty rows
                    if all(cell is None for cell in row):
                        continue

                    # Format as markdown table row
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    md_lines.append(f"| {row_text} |")

            # Save markdown
            md_path = self.output_dir / f"{input_path.stem}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_lines))

            # Excel files typically don't have legal citations
            return ConversionResult(
                md_path=str(md_path),
                citations_found={'pages': [1]},  # Treat as single page
                needs_reconstruction=False,
                doc_type=DocumentType.EXHIBIT
            )

        except Exception as e:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=[f"Excel parsing error: {str(e)}"]
            )

    def email_parser(self, input_path: Path) -> ConversionResult:
        """
        Convert email files (EML, MSG) to markdown.

        Strategy: Extract headers, body, and attachment list.
        """
        if input_path.suffix.lower() == '.msg':
            return self._parse_msg(input_path)
        else:
            return self._parse_eml(input_path)

    def _parse_eml(self, input_path: Path) -> ConversionResult:
        """Parse EML (RFC 822) email files."""
        import email
        from email import policy

        try:
            with open(input_path, 'rb') as f:
                msg = email.message_from_binary_file(f, policy=policy.default)

            md_lines = [
                f"# Email: {msg.get('Subject', '(No Subject)')}\n",
                f"**From:** {msg.get('From', 'Unknown')}",
                f"**To:** {msg.get('To', 'Unknown')}",
                f"**Date:** {msg.get('Date', 'Unknown')}",
                f"**Subject:** {msg.get('Subject', '(No Subject)')}",
                "\n---\n"
            ]

            # Extract body
            body = msg.get_body(preferencelist=('plain', 'html'))
            if body:
                content = body.get_content()
                md_lines.append(content)

            # List attachments
            attachments = [part.get_filename() for part in msg.iter_attachments() if part.get_filename()]
            if attachments:
                md_lines.append("\n**Attachments:**")
                for filename in attachments:
                    md_lines.append(f"- {filename}")

            # Save markdown
            md_path = self.output_dir / f"{input_path.stem}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_lines))

            return ConversionResult(
                md_path=str(md_path),
                citations_found={'pages': [1]},
                needs_reconstruction=False,
                doc_type=DocumentType.EXHIBIT
            )

        except Exception as e:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=[f"Email parsing error: {str(e)}"]
            )

    def _parse_msg(self, input_path: Path) -> ConversionResult:
        """Parse MSG (Outlook) email files."""
        try:
            import extract_msg
        except ImportError:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=["extract-msg not installed. Run: pip install extract-msg"]
            )

        try:
            msg = extract_msg.Message(input_path)

            md_lines = [
                f"# Email: {msg.subject or '(No Subject)'}\n",
                f"**From:** {msg.sender or 'Unknown'}",
                f"**To:** {msg.to or 'Unknown'}",
                f"**Date:** {msg.date or 'Unknown'}",
                f"**Subject:** {msg.subject or '(No Subject)'}",
                "\n---\n",
                msg.body or "(No body text)"
            ]

            # List attachments
            if msg.attachments:
                md_lines.append("\n**Attachments:**")
                for attachment in msg.attachments:
                    md_lines.append(f"- {attachment.longFilename or attachment.shortFilename}")

            msg.close()

            # Save markdown
            md_path = self.output_dir / f"{input_path.stem}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_lines))

            return ConversionResult(
                md_path=str(md_path),
                citations_found={'pages': [1]},
                needs_reconstruction=False,
                doc_type=DocumentType.EXHIBIT
            )

        except Exception as e:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=[f"MSG parsing error: {str(e)}"]
            )

    def powerpoint_parser(self, input_path: Path) -> ConversionResult:
        """
        Convert PowerPoint files to markdown.

        Strategy: Extract text from each slide, preserve slide structure.
        """
        try:
            from pptx import Presentation
        except ImportError:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=["python-pptx not installed. Run: pip install python-pptx"]
            )

        try:
            prs = Presentation(input_path)
            md_lines = [f"# {input_path.stem}\n"]

            for slide_num, slide in enumerate(prs.slides, 1):
                md_lines.append(f"\n## Slide {slide_num}\n")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        md_lines.append(shape.text)

            # Save markdown
            md_path = self.output_dir / f"{input_path.stem}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_lines))

            return ConversionResult(
                md_path=str(md_path),
                citations_found={'pages': list(range(1, len(prs.slides) + 1))},
                needs_reconstruction=False,
                doc_type=DocumentType.EXHIBIT
            )

        except Exception as e:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=[f"PowerPoint parsing error: {str(e)}"]
            )

    def text_parser(self, input_path: Path) -> ConversionResult:
        """
        Handle plain text and markdown files.

        Strategy: Copy to output directory, treat as single page.
        """
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Save to output directory
            md_path = self.output_dir / f"{input_path.stem}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(content)

            return ConversionResult(
                md_path=str(md_path),
                citations_found={'pages': [1]},
                needs_reconstruction=False,
                doc_type=DocumentType.UNKNOWN
            )

        except Exception as e:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=[f"Text parsing error: {str(e)}"]
            )

    def textract_fallback(self, input_path: Path) -> ConversionResult:
        """
        Fallback for unsupported formats using textract.

        Strategy: Use textract universal extractor.
        """
        try:
            import textract
        except ImportError:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=["textract not installed. Run: pip install textract"]
            )

        try:
            text = textract.process(str(input_path)).decode('utf-8')

            # Save as markdown
            md_path = self.output_dir / f"{input_path.stem}.md"
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(f"# {input_path.stem}\n\n{text}")

            return ConversionResult(
                md_path=str(md_path),
                citations_found={'pages': [1]},
                needs_reconstruction=True,
                doc_type=DocumentType.UNKNOWN
            )

        except Exception as e:
            return ConversionResult(
                md_path="",
                citations_found={},
                needs_reconstruction=True,
                errors=[f"Textract error: {str(e)}"]
            )


def main():
    """Test format handler on a sample file."""
    import sys
    if len(sys.argv) < 3:
        print("Usage: python format_handlers.py <input_file> <output_dir>")
        sys.exit(1)

    handler = FormatHandler(sys.argv[2])
    result = handler.convert(Path(sys.argv[1]))

    print(result.citation_coverage_summary())


if __name__ == "__main__":
    main()
